import fitz  # PyMuPDF
from pptx import Presentation
import openpyxl


def extract_text(file_path: str, file_type: str) -> str:
    if file_type == "pdf":
        return _extract_pdf(file_path)
    elif file_type == "pptx":
        return _extract_pptx(file_path)
    elif file_type in ("xlsx", "xls"):
        return _extract_excel(file_path)
    else:
        return ""


def _extract_pdf(file_path: str) -> str:
    doc = fitz.open(file_path)
    pages = []
    for page in doc:
        text = page.get_text()
        if text.strip():
            pages.append(text)
    doc.close()
    return "\n\n---\n\n".join(pages)


def _extract_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        texts.append(line)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        texts.append(row_text)
        if texts:
            slides.append(f"[슬라이드 {i}]\n" + "\n".join(texts))
    return "\n\n---\n\n".join(slides)


def _extract_excel(file_path: str) -> str:
    wb = openpyxl.load_workbook(file_path, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            vals = [str(c) if c is not None else "" for c in row]
            line = " | ".join(vals)
            if line.strip(" |"):
                rows.append(line)
        if rows:
            sheets.append(f"[시트: {sheet_name}]\n" + "\n".join(rows))
    wb.close()
    return "\n\n---\n\n".join(sheets)
